#!/usr/bin/env python3
"""Create Kobold Metals presentation using python-pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Add blank slide
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# Set background color (dark blue)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(10, 22, 40)

# Title
title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.4), Inches(9.2), Inches(0.8))
title_frame = title_box.text_frame
title_frame.word_wrap = True
p = title_frame.paragraphs[0]
p.text = "Kobold Metals - AI-Driven Mineral Exploration Revolution"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(96, 165, 250)

# Subtitle
subtitle_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.1), Inches(9.2), Inches(0.4))
subtitle_frame = subtitle_box.text_frame
p = subtitle_frame.paragraphs[0]
p.text = "Redefining Critical Mineral Discovery with Machine Learning"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(148, 163, 184)

# Left column - Technology Core
tech_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.8), Inches(3.8), Inches(1.5))
tech_frame = tech_box.text_frame
tech_frame.word_wrap = True

p = tech_frame.paragraphs[0]
p.text = "🤖 Technology Core"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = RGBColor(251, 191, 36)
p.space_after = Pt(8)

bullets = [
    "Machine Learning + Big Data driven exploration",
    "Integrating geology, geochemistry, geophysics, and remote sensing",
    "Predicting mineralization zones, improving success rate"
]

for bullet in bullets:
    p = tech_frame.add_paragraph()
    p.text = bullet
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(203, 213, 225)
    p.level = 0
    p.space_after = Pt(6)

# Left column - Investment
invest_box = slide.shapes.add_textbox(Inches(0.4), Inches(3.5), Inches(3.8), Inches(1.2))
invest_frame = invest_box.text_frame
invest_frame.word_wrap = True

p = invest_frame.paragraphs[0]
p.text = "💰 Investment Backing"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = RGBColor(251, 191, 36)
p.space_after = Pt(8)

bullets = [
    "Bill Gates' Breakthrough Energy Ventures",
    "Top Silicon Valley VCs",
    "Valuation: $1B+ (Unicorn status)"
]

for bullet in bullets:
    p = invest_frame.add_paragraph()
    p.text = bullet
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(203, 213, 225)
    p.level = 0
    p.space_after = Pt(6)

# Right column - Timeline
timeline_data = [
    ("2018", "Founded", "Silicon Valley AI team + Geology experts", False),
    ("2022 🎯", "Mingomba Copper Discovery, Zambia", "AI prediction validated • World-class copper deposit • BREAKTHROUGH", True),
    ("2026 🚀", "DRC Lithium Project Launch", "World's largest lithium exploration • Critical battery metals", False)
]

y_pos = 1.8
for year, title, desc, highlight in timeline_data:
    # Year
    year_box = slide.shapes.add_textbox(Inches(4.5), Inches(y_pos), Inches(5), Inches(0.3))
    year_frame = year_box.text_frame
    p = year_frame.paragraphs[0]
    p.text = year
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)

    # Title
    title_box = slide.shapes.add_textbox(Inches(4.5), Inches(y_pos + 0.3), Inches(5), Inches(0.25))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255) if highlight else RGBColor(255, 255, 255)

    # Description
    desc_box = slide.shapes.add_textbox(Inches(4.5), Inches(y_pos + 0.55), Inches(5), Inches(0.4))
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True
    p = desc_frame.paragraphs[0]
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(148, 163, 184)

    # Highlight marker for Mingomba
    if highlight:
        marker = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(4.3), Inches(y_pos), Inches(0.08), Inches(1.0)
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = RGBColor(251, 191, 36)
        marker.line.fill.background()

    y_pos += 1.1

# Footer
footer_box = slide.shapes.add_textbox(Inches(0.4), Inches(5.0), Inches(9.2), Inches(0.3))
footer_frame = footer_box.text_frame
p = footer_frame.paragraphs[0]
p.text = "⚡ Strategic Mission: Powering the global energy transition with critical minerals (Li/Cu/Co/Ni)"
p.font.size = Pt(11)
p.font.color.rgb = RGBColor(100, 116, 139)

# Save
prs.save('Kobold-Metals-Presentation.pptx')
print("✓ Presentation created: Kobold-Metals-Presentation.pptx")
